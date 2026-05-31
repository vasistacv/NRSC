"""
Professional Nature-style PPT — Light theme, clean, minimal.
Slide 1: Title + single-paragraph abstract (statistical downscaling)
Remaining slides: Full-size figures only, no text clutter.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
ACCENT      = RGBColor(0x00, 0x5C, 0x53)  # deep teal
LINE_COLOR  = RGBColor(0xDD, 0xDD, 0xDD)
BG_COLOR    = RGBColor(0xFB, 0xFB, 0xFB)  # near-white


def set_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_line(slide, left, top, width, color=LINE_COLOR, thickness=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.height = int(thickness)
    return shape


def add_text(slide, left, top, width, height, text, size=14,
             color=BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             font="Georgia", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE + ABSTRACT
# ═════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(s1)

# Thin accent bar at very top
add_line(s1, Inches(0), Inches(0), Inches(13.333), color=ACCENT, thickness=Pt(4))

# Title
add_text(s1, Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.9),
         "Statistical Downscaling of Rainfall Using Neural Networks\n"
         "for Station-Level Prediction over Telangana, India",
         size=28, color=BLACK, bold=True, align=PP_ALIGN.LEFT,
         line_spacing=36)

# Thin separator
add_line(s1, Inches(1.0), Inches(2.0), Inches(2.5), color=ACCENT, thickness=Pt(2))

# Abstract — single clean paragraph
abstract = (
    "Skilful prediction of station-level rainfall over the Indian subcontinent remains a "
    "persistent challenge, as global numerical weather prediction models such as ECMWF IFS "
    "and NCEP GFS operate at horizontal resolutions too coarse to resolve mesoscale convective "
    "processes that dominate monsoonal precipitation. In this study, we develop a compact neural "
    "network-based statistical downscaling framework that ingests gridded ECMWF forecasts alongside "
    "physics-derived scalar indices to produce calibrated daily rainfall predictions "
    "at seven automatic weather stations across Telangana during the Indian Summer Monsoon "
    "(June\u2013September, 2015\u20132024). Through systematic variable reduction from 19 to 7 input "
    "channels, the framework achieves improved extreme-event detection while requiring fewer than "
    "72,000 trainable parameters. Evaluated under a rigorous Leave-One-Year-Out Cross-Validation "
    "protocol spanning ten monsoon seasons, the framework consistently outperforms raw ECMWF "
    "forecasts across all categorical skill metrics\u2014achieving higher Critical Success Index, "
    "Symmetric Extremal Dependence Index, and temporal correlation for both general and extreme "
    "rainfall events. These results "
    "demonstrate that lightweight neural network post-processing offers a robust and operationally "
    "viable pathway for bridging the scale gap between global NWP output and the hyper-local "
    "rainfall information demanded by agricultural, hydrological, and disaster-management applications "
    "in data-sparse tropical regions."
)

add_text(s1, Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.5),
         abstract,
         size=15, color=DARK_GRAY, bold=False, align=PP_ALIGN.JUSTIFY,
         font="Georgia", line_spacing=24)

# Bottom accent bar
add_line(s1, Inches(0), Inches(7.46), Inches(13.333), color=ACCENT, thickness=Pt(4))


# ═════════════════════════════════════════════════════════════════════════
# FIGURE SLIDES — images only, clean and full-size
# ═════════════════════════════════════════════════════════════════════════
figure_files = [
    r"D:\NEW_NRSC\paper_figures\Fig1_Study_Area_Map.png",
    r"D:\NEW_NRSC\paper_figures\Fig_Architecture.png",
    r"D:\NEW_NRSC\paper_figures\Fig_Correlation_Heatmap.png",
    r"D:\NEW_NRSC\paper_figures\Fig3_Station_Heatmap.png",
    r"D:\NEW_NRSC\paper_figures\Fig4_Radar_Comparison.png",
    r"D:\NEW_NRSC\paper_figures\Fig5_Split_Comparison.png",
    r"D:\NEW_NRSC\paper_figures\Fig6_Spatial_Performance.png",
    r"D:\NEW_NRSC\paper_figures\Fig7_Station_Distribution.png",
    r"D:\NEW_NRSC\paper_figures\Fig8_LOYO_Yearly.png",
    r"D:\NEW_NRSC\paper_figures\Fig9_Comprehensive_Table.png",
    r"D:\NEW_NRSC\paper_figures\Fig10_Improvement_Summary.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Chevella.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Hayathnagar.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Ibrahimpatnam.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Kondurg.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Maheshwaram.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Saroornagar.png",
    r"D:\NEW_NRSC\paper_figures\Fig11_Timeseries_2019_Yacharam.png",
]

for img_path in figure_files:
    if not os.path.exists(img_path):
        print(f"  SKIP: {img_path}")
        continue

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)

    # Top accent line
    add_line(slide, Inches(0), Inches(0), Inches(13.333), color=ACCENT, thickness=Pt(3))

    # Load image dimensions
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h

    # Fill as much of the slide as possible with padding
    max_w = Inches(12.5)
    max_h = Inches(7.0)

    if aspect > (max_w / max_h):
        w = max_w
        h = int(w / aspect)
    else:
        h = max_h
        w = int(h * aspect)

    left = int((Inches(13.333) - w) / 2)
    top = int((Inches(7.5) - h) / 2)

    slide.shapes.add_picture(img_path, left, top, w, h)

    # Bottom accent line
    add_line(slide, Inches(0), Inches(7.47), Inches(13.333), color=ACCENT, thickness=Pt(3))


# ═════════════════════════════════════════════════════════════════════════
# THANK YOU SLIDE
# ═════════════════════════════════════════════════════════════════════════
s_end = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(s_end)
add_line(s_end, Inches(0), Inches(0), Inches(13.333), color=ACCENT, thickness=Pt(4))

add_text(s_end, Inches(2.0), Inches(2.8), Inches(9.3), Inches(1.0),
         "Thank You",
         size=44, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font="Georgia")

add_line(s_end, Inches(5.5), Inches(4.2), Inches(2.3), color=ACCENT, thickness=Pt(2))

add_text(s_end, Inches(2.0), Inches(4.5), Inches(9.3), Inches(0.6),
         "Because of your guidance, I learned a lot of things. Thank you, Sir.",
         size=18, color=MID_GRAY, bold=False, align=PP_ALIGN.CENTER, font="Georgia")

add_line(s_end, Inches(0), Inches(7.46), Inches(13.333), color=ACCENT, thickness=Pt(4))

# ── Save ─────────────────────────────────────────────────────────────────
out = r"D:\NEW_NRSC\paper_figures\Rainfall_Prediction_Presentation.pptx"
prs.save(out)
print(f"\n[OK] Saved -> {out}")
print(f"     Slides: {len(prs.slides)}")
